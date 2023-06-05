# 将 ppt 的备注导出到 ppt 所在路径且文件名相同的 word 中。
# 教程：https://www.zkxjob.com/4787
from pptx import Presentation  # pip install python-pptx
#创建并写入word文档
import docx  # pip install python-docx

from docx.shared import Cm
from docx.enum.section import WD_ORIENTATION

import os

# 使用时修改这个 ppt 所在路径的变量
ppt_dir = 'D:/doc/doctor/答辩/3_答辩'
file_name = '答辩_王海东'

input_file = os.path.join(ppt_dir, '答辩_王海东.pptx')
ppt = Presentation(input_file)

#创建内存中的word文档对象
doc = docx.Document()

notes = []

for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note.
    textNote = slide.notes_slide.notes_text_frame.text
    notes.append((page, textNote))
    # 写入若干段落
    doc.add_paragraph(str(page+1))
    doc.add_paragraph(textNote)

print(notes)


# 为了便于只能单面打印的打印机打印双面，当时奇数页面的时候，需要在最后插入空白页。
# 参考：https://icode9.com/content-1-961624.html
# 奇数页PDF插入空白页
import pythoncom
import win32com.client as win32
def get_word_pages(path):
    pythoncom.CoInitialize()
    word =win32.Dispatch("Word.Application") # client.
    doc = word.Documents.Open(path)
    word.ActiveDocument.Repaginate()
    pages=word.ActiveDocument.ComputeStatistics(2)
    print(pages)
    doc.Close()
    return pages
p=get_word_pages(path)

if len(len(doc.pages)) % 2 == 1:
    doc.add_page_break()

# 设置页边距
doc.sections[0].top_margin = Cm(0)  # 3.7
doc.sections[0].bottom_margin = Cm(0)  # 3.5
doc.sections[0].left_margin = Cm(0.2)  # 2.8
doc.sections[0].right_margin = Cm(0.2)  # 2.6
# 保存
doc_path = os.path.join(ppt_dir, "答辩_王海东.docx")
doc.save(doc_path)

# pip install docx2pdf
# pip install pywin32==223
from docx2pdf import convert
convert(doc_path, os.path.join(ppt_dir, file_name + '.pdf'))


# pathofcwd = input("请输入要处理的PDF的文件路径")